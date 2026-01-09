
from ragSystems.ragProcessor import ragProcessor
from langchain_text_splitters import RecursiveCharacterTextSplitter
from services.prompt.promptProcessor import taskProcessor
import pdfplumber
import threading
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from PIL import Image
from docx import Document
from pypdf import PdfReader
import os
from io import BytesIO
import zipfile
import re
import unicodedata
# For better PDF image extraction, import PyMuPDF (fitz)
import fitz  # pip install PyMuPDF

# Assuming ClipSearchService is imported from the appropriate module; adjust the import path as needed
from ragSystems.imageRag import \
    ClipSearchService  # Replace 'your_module' with the actual module path containing ClipSearchService


class documentProcessor:
    def __init__(self, collection_name):
        self.taskHandler = taskProcessor()
        self.ragProcessor = ragProcessor(collection_name=collection_name)
        self.image_service = ClipSearchService(collection_name=str(collection_name) + "_img")
        self.BASE_IMAGE_DIR = r"D:\scoriX_agent\data\referenceImages"
        os.makedirs(self.BASE_IMAGE_DIR, exist_ok=True)

    def sanitize_path(self, path: str) -> str:
        """
        Remove or replace invalid filename characters to make paths safe for Windows.
        Strips control characters like \x07 and normalizes.
        """
        if not path or not isinstance(path, str):
            return path

        # Normalize unicode (e.g., NFKD to handle composed chars)
        path = unicodedata.normalize('NFKD', path)

        # Remove control characters (ASCII 0-31 except tab/newline)
        path = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', path)

        # Replace other invalid Windows chars with '_'
        invalid_chars = r'[<>:"/\\|?*]'
        path = re.sub(invalid_chars, '_', path)

        # Trim excessive whitespace/dots
        path = re.sub(r'\s+', ' ', path.strip())
        path = re.sub(r'\.+', '.', path)

        return path

    def chunkTextPerPage(self, page_data, chunkSize=600, chunkOverlap=150):
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunkSize,
            chunk_overlap=chunkOverlap,
            separators=["\n\n", "\n", ".", " ", ""]
        )

        allChunks = []
        for page_text in page_data:
            if page_text:
                textChunks = splitter.split_text(page_text)
                allChunks.extend([f"{chunk}" for chunk in textChunks])
        return allChunks

    def process_with_llm(self, content):
        """Function to handle LLM processing synchronously"""
        try:
            response = self.taskHandler.tableSummaryProcessor(content=content)
            return response
        except Exception as e:
            print(f"❌ LLM processing failed: {e}")
            return None

    def format_table_for_gemini(self, table):
        """
        Convert extracted table data to Gemini Pro optimized format
        """
        if not table or len(table) < 2:
            return ""

        # Get headers
        headers = [str(cell).strip() if cell else "" for cell in table[0]]

        # Create markdown table
        markdown = "| " + " | ".join(headers) + " |\n"
        markdown += "|" + "|".join([" --- " for _ in headers]) + "|\n"

        # Add data rows
        for row in table[1:]:
            row_data = [str(cell).strip() if cell else "" for cell in row]
            markdown += "| " + " | ".join(row_data) + " |\n"

        return markdown + "\n"

    def extract_images_from_pdf_with_fitz(self, pdf_path):
        """
        Extract images from PDF using PyMuPDF (fitz) for better handling of embedded images.
        Returns list of (save_path, metadata) tuples.
        """
        image_paths = []
        metadatas = []
        try:
            doc = fitz.open(pdf_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images(full=True)
                for img_idx, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n - pix.alpha < 4:  # Skip if not a valid color image
                            # Convert to RGB if necessary
                            if pix.n == 4:  # CMYK
                                pix = fitz.utils.get_colorspace_rgb().to_pixmap(pix)
                            elif pix.n == 1:  # Grayscale, convert to RGB
                                mat = fitz.Matrix(2.0, 2.0)  # Optional: upscale
                                pix = fitz.get_pixmap(matrix=mat, colorspace=fitz.csRGB, clip=None)
                            # Determine extension
                            ext = 'png' if pix.alpha else 'jpg'
                            base_name = os.path.splitext(os.path.basename(pdf_path))[0]
                            save_path = os.path.join(self.BASE_IMAGE_DIR,
                                                     f"{base_name}_page_{page_num + 1}_img_{img_idx}.{ext}")
                            pix.save(save_path)
                            image_paths.append(save_path)
                            metadatas.append({
                                "source_file": pdf_path,
                                "page": page_num + 1,
                                "type": "pdf_image",
                                "width": pix.width,
                                "height": pix.height
                            })
                        pix = None  # Free resources
                    except Exception as e:
                        print(f"Error extracting image {img_idx} from PDF page {page_num + 1} with fitz: {e}")
            doc.close()
        except Exception as e:
            print(f"❌ Error opening PDF with fitz: {e}")
        return image_paths, metadatas

    def extract_from_pdf(self, pdf_path):
        page_data = []
        image_paths = []
        metadatas = []
        text = ""  # Accumulated table content
        flag = False

        try:
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""
                    page_tables = page.extract_tables() or []

                    page_data.append(page_text)

                    current_has_tables = bool(page_tables)
                    if current_has_tables:
                        if not flag:
                            text = ""
                        table_str = ""
                        for i, table in enumerate(page_tables):
                            table_str += f"\n--- TABLE {i + 1} (Page {page_num + 1}) ---\n"
                            formatted_table = self.format_table_for_gemini(table)
                            table_str += formatted_table
                            table_str += "\n" + "=" * 50 + "\n"
                        text += table_str
                        flag = True
                    else:
                        if flag:
                            summary = self.process_with_llm(text)
                            if summary:
                                page_data.append(summary)
                            text = ""
                            flag = False

                # Handle trailing tables
                if flag:
                    summary = self.process_with_llm(text)
                    if summary:
                        page_data.append(summary)

            # Extract images using fitz for better reliability
            pdf_image_paths, pdf_metadatas = self.extract_images_from_pdf_with_fitz(pdf_path)
            image_paths.extend(pdf_image_paths)
            metadatas.extend(pdf_metadatas)

            return [p for p in page_data if p], image_paths, metadatas

        except Exception as e:
            print(f"❌ Error processing PDF file: {e}")
            # Fallback: try to extract images anyway
            _, fallback_images, fallback_meta = self.extract_from_pdf_fallback(pdf_path)
            return [], fallback_images, fallback_meta

    def extract_from_pdf_fallback(self, pdf_path):
        """Fallback extraction using pdfplumber, but with better validation."""
        image_paths = []
        metadatas = []
        try:
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    page_images = page.images or []
                    for img_idx, img in enumerate(page_images):
                        try:
                            image_bytes = img["stream"].get_data()
                            # Basic validation: check if it's a plausible image (e.g., starts with JPEG/PNG magic bytes)
                            if len(image_bytes) < 100:
                                continue  # Too small, likely not an image
                            magic_bytes = image_bytes[:8]
                            if not (magic_bytes.startswith(b'\xFF\xD8\xFF') or  # JPEG
                                    magic_bytes.startswith(b'\x89PNG\r\n\x1a\n') or  # PNG
                                    magic_bytes.startswith(b'GIF87a') or magic_bytes.startswith(b'GIF89a') or  # GIF
                                    magic_bytes.startswith(b'RIFF') and b'WEBP' in image_bytes[:12]):  # WEBP
                                print(f"Skipping non-image stream on page {page_num + 1}")
                                continue
                            img_pil = Image.open(BytesIO(image_bytes))
                            ext = img_pil.format.lower() if img_pil.format else 'png'
                            base_name = os.path.splitext(os.path.basename(pdf_path))[0]
                            save_path = os.path.join(self.BASE_IMAGE_DIR,
                                                     f"{base_name}_page_{page_num + 1}_img_{img_idx}_fallback.{ext}")
                            img_pil.save(save_path)
                            image_paths.append(save_path)
                            metadatas.append({
                                "source_file": pdf_path,
                                "page": page_num + 1,
                                "type": "pdf_image_fallback"
                            })
                        except Exception as e:
                            print(f"Error in fallback extraction for PDF page {page_num + 1}: {e}")
        except Exception as e:
            print(f"❌ Fallback PDF processing failed: {e}")
        return [], image_paths, metadatas

    def extractTextAndTablesDocx(self, docxPath):
        """Extract text and tables from DOCX files"""
        page_data = []
        image_paths = []
        metadatas = []
        try:
            if not docxPath.lower().endswith('.docx'):
                raise ValueError("Provided file is not a .docx Word document.")

            doc = Document(docxPath)

            # Extract text content
            full_text_parts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
            full_text = "\n".join(full_text_parts)

            page_data = [full_text] if full_text else []

            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if len(table_rows) > 1 and any(any(cell for cell in row) for row in table_rows[1:]):
                    table_str = f"\n--- TABLE {table_idx + 1} ---\n"
                    formatted_table = self.format_table_for_gemini(table_rows)
                    table_str += formatted_table
                    summary = self.process_with_llm(table_str)
                    if summary:
                        page_data.append(summary)

            # Extract images
            with zipfile.ZipFile(docxPath) as docx_zip:
                for file_name in docx_zip.namelist():
                    if file_name.startswith('word/media/') and not file_name.endswith('/'):
                        try:
                            image_data = docx_zip.read(file_name)
                            # Validate image
                            if len(image_data) < 100:
                                continue
                            magic_bytes = image_data[:8]
                            if not (magic_bytes.startswith(b'\xFF\xD8\xFF') or
                                    magic_bytes.startswith(b'\x89PNG\r\n\x1a\n') or
                                    magic_bytes.startswith(b'GIF87a') or magic_bytes.startswith(b'GIF89a')):
                                print(f"Skipping non-image in DOCX: {file_name}")
                                continue
                            img_pil = Image.open(BytesIO(image_data))
                            ext = img_pil.format.lower() if img_pil.format else 'png'
                            base_name = os.path.splitext(os.path.basename(docxPath))[0]
                            save_path = os.path.join(self.BASE_IMAGE_DIR, f"{base_name}_img_{len(image_paths)}.{ext}")
                            img_pil.save(save_path)
                            image_paths.append(save_path)
                            metadatas.append({
                                "source_file": docxPath,
                                "type": "docx_image",
                                "image_id": len(image_paths)
                            })
                        except Exception as e:
                            print(f"Error extracting image {file_name} from DOCX: {e}")

            return [p for p in page_data if p], image_paths, metadatas

        except Exception as e:
            print(f"❌ Error processing DOCX file: {e}")
            return [], [], []

    def extract_from_pptx(self, pptx_path):
        """Extract text and tables from PPTX files"""
        page_data = []
        image_paths = []
        metadatas = []
        try:
            if not pptx_path.lower().endswith('.pptx'):
                raise ValueError("Provided file is not a .pptx PowerPoint file.")

            prs = Presentation(pptx_path)

            for slide_num, slide in enumerate(prs.slides):
                slide_texts = []
                slide_images = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            slide_texts.append(text)
                    elif shape.has_table:
                        table = shape.table
                        table_rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                        if len(table_rows) > 1 and any(any(cell for cell in row) for row in table_rows[1:]):
                            table_str = f"\n--- TABLE (Slide {slide_num + 1}) ---\n"
                            formatted_table = self.format_table_for_gemini(table_rows)
                            table_str += formatted_table
                            summary = self.process_with_llm(table_str)
                            if summary:
                                slide_texts.append(summary)
                    # Extract images (picture shapes)
                    elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            image = shape.image
                            image_data = image.blob
                            # Validate
                            if len(image_data) < 100:
                                continue
                            magic_bytes = image_data[:8]
                            if not (magic_bytes.startswith(b'\xFF\xD8\xFF') or
                                    magic_bytes.startswith(b'\x89PNG\r\n\x1a\n') or
                                    magic_bytes.startswith(b'GIF87a') or magic_bytes.startswith(b'GIF89a')):
                                print(f"Skipping non-image in PPTX slide {slide_num + 1}")
                                continue
                            img_pil = Image.open(BytesIO(image_data))
                            ext = img_pil.format.lower() if img_pil.format else 'png'
                            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
                            save_path = os.path.join(self.BASE_IMAGE_DIR,
                                                     f"{base_name}_slide_{slide_num + 1}_img_{len(slide_images)}.{ext}")
                            img_pil.save(save_path)
                            image_paths.append(save_path)
                            metadatas.append({
                                "source_file": pptx_path,
                                "slide": slide_num + 1,
                                "type": "pptx_image"
                            })
                            slide_images.append("image_placeholder")  # Placeholder to indicate image presence
                        except Exception as e:
                            print(f"Error extracting image from PPTX slide {slide_num + 1}: {e}")

                if slide_texts or slide_images:
                    slide_text = "\n".join(slide_texts)
                    page_data.append(slide_text)

            return [p for p in page_data if p], image_paths, metadatas

        except Exception as e:
            print(f"❌ Error processing PPTX file: {e}")
            return [], [], []

    def extract_from_txt(self, txt_path):
        """Extract text from TXT files"""
        try:
            if not txt_path.lower().endswith('.txt'):
                raise ValueError("Provided file is not a .txt file.")

            with open(txt_path, 'r', encoding='utf-8') as f:
                text = f.read().strip()

            page_data = [text] if text else []
            return page_data, [], []

        except Exception as e:
            print(f"❌ Error processing TXT file: {e}")
            return [], [], []

    def escape_windows_path(self,path: str) -> str:
        return path.replace('\\', '\\\\')
    def processor(self, pdf_path):
        if not pdf_path:
            return

        # Sanitize the path before processing
        pdf_path = self.escape_windows_path(pdf_path)
        print(f"🔍 Sanitized path: {pdf_path}")  # Log for debugging

        if not os.path.exists(pdf_path):
            print(f"❌ Sanitized path does not exist: {pdf_path}")
            return

        ext = pdf_path.lower().split('.')[-1] if '.' in pdf_path else ''
        if ext == 'pdf':
            page_data, image_paths, metadatas = self.extract_from_pdf(pdf_path)
            print("image_paths")
            print(image_paths)
        elif ext == 'docx':
            page_data, image_paths, metadatas = self.extractTextAndTablesDocx(pdf_path)
        elif ext == 'pptx':
            page_data, image_paths, metadatas = self.extract_from_pptx(pdf_path)
        elif ext == 'txt':
            page_data, image_paths, metadatas = self.extract_from_txt(pdf_path)
        else:
            print(f"❌ Unsupported file type: {ext}")
            return

        if not page_data:
            print("❌ No content extracted from the file.")
            return

        chunks = self.chunkTextPerPage(page_data)

        if chunks:
            self.ragProcessor.docStoring(chunks)
            print(f"✅ Processed and stored {len(chunks)} chunks in RAG collection.")

        if image_paths:
            self.image_service.index_images(image_paths, metadatas)
            print(f"✅ Processed and stored {len(image_paths)} images in image collection.")
        else:
            print("No images found to process.")

    def search(self, query):
        return self.ragProcessor.search(query)#       chunks  = [
#     "Artificial Intelligence (AI) is a branch of computer science that aims to create machines capable of intelligent behavior.",
#     "Machine learning is a subset of AI that allows systems to learn patterns from data and improve performance over time without explicit programming.",
#     "Deep learning is a type of machine learning that uses neural networks with many layers to model complex patterns in data.",
#     "Natural Language Processing (NLP) enables machines to understand, interpret, and generate human language.",
#     "Computer vision is a field of AI that trains computers to interpret and process visual information from the world.",
#     "Reinforcement learning is a type of machine learning where agents learn to make decisions by receiving rewards or penalties.",
#     "Supervised learning uses labeled datasets to train models to predict outputs from inputs.",
#     "Unsupervised learning identifies hidden patterns in data without pre-existing labels.",
#     "Semi-supervised learning uses a combination of labeled and unlabeled data to improve learning accuracy.",
#     "Transfer learning leverages knowledge from one model trained on a task to improve performance on a different but related task.",
#     "Convolutional Neural Networks (CNNs) are widely used in image recognition and processing tasks.",
#     "Recurrent Neural Networks (RNNs) are designed for sequential data, such as time series or language modeling.",
#     "Generative Adversarial Networks (GANs) consist of two networks, a generator and a discriminator, competing to produce realistic outputs.",
#     "AI ethics concerns the moral implications and responsible use of artificial intelligence technologies.",
#     "Bias in AI models can arise from skewed training data, leading to unfair or inaccurate predictions.",
#     "Explainable AI (XAI) aims to make AI decisions understandable and transparent to humans.",
#     "Data preprocessing is essential to clean and prepare raw data before feeding it into machine learning models.",
#     "Feature engineering involves creating meaningful input variables that improve model performance.",
#     "Hyperparameter tuning is the process of optimizing model parameters that are not learned during training.",
#     "Overfitting occurs when a model learns training data too well, failing to generalize to unseen data.",
#     "Regularization techniques like L1 and L2 are used to prevent overfitting in machine learning models.",
#     "Gradient descent is an optimization algorithm used to minimize the loss function in neural networks.",
#     "Backpropagation is the method of calculating gradients for all weights in a neural network to update them efficiently.",
#     "Activation functions, such as ReLU or sigmoid, introduce non-linearity into neural networks to model complex patterns.",
#     "Loss functions quantify the difference between predicted and actual outputs, guiding model training.",
#     "Optimizers like Adam or SGD adjust model weights to reduce the loss function during training.",
#     "Cross-validation is a technique to evaluate a model's performance on multiple data splits.",
#     "Precision, recall, and F1-score are common metrics to evaluate classification models.",
#     "Clustering algorithms, like K-Means, group similar data points without labeled outputs.",
#     "Dimensionality reduction techniques, like PCA, reduce the number of features while preserving important information."
# ]